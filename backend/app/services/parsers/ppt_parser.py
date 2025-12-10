"""
Advanced PPT Parser with shape and table extraction
Adapted from RAGFlow PptParser
"""
import logging
from io import BytesIO
from typing import List

from pptx import Presentation

from .utils import num_tokens_from_string

logger = logging.getLogger(__name__)


class AdvancedPptParser:
    """
    Advanced PowerPoint parser extracting text, tables, and shapes.
    """
    
    def __init__(self):
        self.presentation = None
    
    def parse(self, file_path: str = None, binary: bytes = None) -> str:
        """
        Parse PowerPoint and return full text.
        
        Args:
            file_path: Path to PPT/PPTX file
            binary: Binary content
            
        Returns:
            Full extracted text
        """
        # Load presentation
        if binary:
            self.presentation = Presentation(BytesIO(binary))
        elif file_path:
            self.presentation = Presentation(file_path)
        else:
            raise ValueError("Either file_path or binary must be provided")
        
        # Extract from all slides
        slides_text = []
        
        for slide_num, slide in enumerate(self.presentation.slides, 1):
            slide_content = self._extract_slide(slide, slide_num)
            if slide_content:
                slides_text.append(f"=== Slide {slide_num} ===\n{slide_content}")
        
        return "\n\n".join(slides_text)
    
    def _extract_slide(self, slide, slide_num: int) -> str:
        """
        Extract content from a single slide.
        
        Args:
            slide: python-pptx Slide object
            slide_num: Slide number
            
        Returns:
            Slide content
        """
        parts = []
        
        # Extract from all shapes
        for shape in slide.shapes:
            content = self._extract_shape(shape)
            if content:
                parts.append(content)
        
        return "\n\n".join(parts)
    
    def _extract_shape(self, shape) -> str:
        """
        Extract content from a shape.
        
        Args:
            shape: python-pptx Shape object
            
        Returns:
            Shape content
        """
        content_parts = []
        
        # Text frame
        if hasattr(shape, "text") and shape.text:
            text = shape.text.strip()
            if text:
                content_parts.append(text)
        
        # Table
        if hasattr(shape, "table"):
            table_text = self._extract_table(shape.table)
            if table_text:
                content_parts.append(table_text)
        
        # Group shapes (nested)
        if hasattr(shape, "shapes"):
            for sub_shape in shape.shapes:
                sub_content = self._extract_shape(sub_shape)
                if sub_content:
                    content_parts.append(sub_content)
        
        return "\n".join(content_parts)
    
    def _extract_table(self, table) -> str:
        """
        Extract and format table content.
        
        Args:
            table: python-pptx Table object
            
        Returns:
            Formatted table text
        """
        if not table or len(table.rows) < 2:
            return ""
        
        # Extract rows
        rows = []
        for row in table.rows:
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip()
                cells.append(cell_text)
            rows.append(cells)
        
        # Use first row as headers
        headers = rows[0]
        data_rows = rows[1:]
        
        # Format rows
        formatted_rows = []
        for row in data_rows:
            pairs = []
            for i, value in enumerate(row):
                if value:
                    if i < len(headers) and headers[i]:
                        pairs.append(f"{headers[i]}: {value}")
                    else:
                        pairs.append(value)
            
            if pairs:
                formatted_rows.append("; ".join(pairs))
        
        return "\n".join(formatted_rows)
    
    def parse_with_chunks(self, file_path: str = None, binary: bytes = None,
                          chunk_token_count: int = 128,
                          delimiter: str = "\n") -> List[str]:
        """
        Parse PowerPoint with intelligent chunking.
        
        Args:
            file_path: Path to PPT/PPTX
            binary: Binary content
            chunk_token_count: Target tokens per chunk
            delimiter: Delimiter for chunking
            
        Returns:
            List of text chunks
        """
        # Get full text
        full_text = self.parse(file_path, binary)
        
        # Use chunking logic from txt_parser
        from .txt_parser import AdvancedTxtParser
        
        txt_parser = AdvancedTxtParser()
        return txt_parser._chunk_text(full_text, chunk_token_count, delimiter)
    
    def __call__(self, file_path: str, binary: bytes = None) -> str:
        """
        Callable interface returning text.
        
        Returns:
            Extracted text
        """
        return self.parse(file_path, binary)
